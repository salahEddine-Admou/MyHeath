#!/usr/bin/env python3
"""Genere RAPPORT_PFE_HERACARE.pdf et PRESENTATION_PFE_HERACARE.pptx"""

from pathlib import Path

from fpdf import FPDF
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUT = Path(__file__).resolve().parent.parent / "docs"
OUT.mkdir(exist_ok=True)

ROSE = RGBColor(0xC9, 0x2D, 0x55)
INK = RGBColor(0x1A, 0x12, 0x16)
MUTED = RGBColor(0x6B, 0x55, 0x60)


def T(s: str) -> str:
    repl = {
        "é": "e", "è": "e", "ê": "e", "ë": "e", "à": "a", "â": "a",
        "ù": "u", "û": "u", "ô": "o", "ö": "o", "î": "i", "ï": "i",
        "ç": "c", "É": "E", "È": "E", "À": "A", "Ù": "U", "Ô": "O",
        "Î": "I", "Ç": "C", "'": "'", "'": "'", '"': '"', '"': '"',
        "–": "-", "—": "-", "…": "...", "«": '"', "»": '"',
        "≤": "<=", "≥": ">=", "→": "->", "↔": "<->", "×": "x", "·": "-",
    }
    for a, b in repl.items():
        s = s.replace(a, b)
    return s.encode("latin-1", "replace").decode("latin-1")


class RapportPDF(FPDF):
    def header(self):
        if self.page_no() == 1:
            return
        self.set_font("Helvetica", "I", 9)
        self.set_text_color(168, 33, 69)
        self.cell(0, 8, T("HeraCare - Rapport de PFE"), align="L")
        self.ln(4)
        self.set_draw_color(201, 45, 85)
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(6)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(120, 120, 120)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

    def title_page(self):
        self.add_page()
        self.ln(30)
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(100, 100, 100)
        self.multi_cell(0, 8, T("ROYAUME DU MAROC\nProjet de Fin d'Etudes - Sciences de l'Ingenieur"), align="C")
        self.ln(20)
        self.set_font("Helvetica", "B", 22)
        self.set_text_color(201, 45, 85)
        self.multi_cell(0, 12, "HeraCare", align="C")
        self.ln(4)
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(26, 18, 22)
        self.multi_cell(
            0, 8,
            T("Conception et realisation d'une plateforme intelligente\nde telemedecine et de suivi de sante feminine"),
            align="C",
        )
        self.ln(16)
        self.set_font("Helvetica", "", 11)
        self.set_text_color(80, 80, 80)
        self.multi_cell(
            0, 7,
            T("Stack MERN - Docker - AES-256 - Socket.io - Claude AI\nFiliere : Genie Informatique / Genie Logiciel"),
            align="C",
        )
        self.ln(30)
        self.set_font("Helvetica", "", 11)
        self.multi_cell(
            0, 7,
            T("Auteur : Nezha Fekoussa\nAnnee universitaire : 2025-2026"),
            align="C",
        )

    def body(self, text):
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 40, 40)
        self.multi_cell(0, 5.5, T(text))
        self.ln(1)

    def bullet(self, text):
        self.set_x(self.l_margin)
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 40, 40)
        self.multi_cell(0, 5.5, T(f"  - {text}"))
        self.set_x(self.l_margin)

    def h1(self, text):
        self.set_x(self.l_margin)
        self.ln(4)
        self.set_font("Helvetica", "B", 16)
        self.set_text_color(201, 45, 85)
        self.multi_cell(0, 9, T(text))
        self.ln(2)

    def h2(self, text):
        self.set_x(self.l_margin)
        self.ln(3)
        self.set_font("Helvetica", "B", 13)
        self.set_text_color(26, 18, 22)
        self.multi_cell(0, 7, T(text))
        self.ln(1)


def build_pdf():
    pdf = RapportPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(True, 18)
    pdf.title_page()

    pdf.add_page()
    pdf.h1("Dedicaces & Remerciements")
    pdf.body(
        "Je dedie ce travail a ma famille pour son soutien constant. "
        "Mes remerciements s'adressent a mon encadrant, au corps professoral, "
        "ainsi qu'a toutes les personnes qui ont contribue a la reussite de ce projet."
    )

    pdf.h1("Resume (Francais)")
    pdf.body(
        "HeraCare est une application web de telemedecine dediee a la sante feminine (FemTech). "
        "Elle combine un suivi menstruel predictif, un dossier medical partage securise, "
        "une messagerie temps reel Patient-Medecin, et un assistant IA (Hera AI / Claude). "
        "L'architecture repose sur la stack MERN, conteneurisee via Docker, avec un chiffrement "
        "AES-256-CBC des donnees de sante au repos, conforme aux principes de la loi 09-08 (Maroc)."
    )
    pdf.body("Mots-cles : FemTech, telemedecine, MERN, AES-256, Socket.io, IA Claude, Loi 09-08, PFE")

    pdf.h2("Abstract (English)")
    pdf.body(
        "HeraCare is a FemTech telemedicine platform with predictive cycle tracking, "
        "AES-256 encrypted medical records, real-time messaging, and a Claude-powered health assistant. "
        "Built on the MERN stack with Docker orchestration."
    )

    pdf.add_page()
    pdf.h1("Introduction generale")
    pdf.h2("Contexte")
    pdf.body(
        "La sante des femmes reste insuffisamment adressee par les systemes numeriques, "
        "surtout dans les zones a faible densite medicale. Le marche FemTech croit fortement, "
        "mais l'offre nationale demeure fragmentee. HeraCare propose un suivi algorithmique "
        "du cycle et une consultation a distance securisee, enrichie par l'intelligence artificielle."
    )
    pdf.h2("Problematique")
    pdf.body(
        "Comment concevoir une plateforme de telemedecine feminine qui (1) protege les donnees "
        "de sante sensibles, (2) apporte une valeur clinique via l'analyse predictive et l'IA, "
        "et (3) s'appuie sur une architecture logicielle modulaire digne d'un projet d'ingenierie ?"
    )
    pdf.h2("Objectifs")
    for o in [
        "Architecture MERN modulaire (Auth, Dossier, Consultation, Engine, IA).",
        "Chiffrement AES-256 des donnees medicales et messages au repos.",
        "Moteur de prediction de cycle et detection d'anomalies (signaux SOPK / endometriose).",
        "Assistant IA Hera (Claude) : chat, journal NL, brief medecin, plan bien-etre.",
        "Deploiement frontend Vercel + backend cloud + MongoDB Atlas.",
    ]:
        pdf.bullet(o)

    pdf.add_page()
    pdf.h1("Chapitre 1 - Etude prealable et besoins")
    pdf.h2("Etude de l'existant")
    pdf.body(
        "Les applications internationales (Clue, Flo, Ada) offrent un suivi de cycle grand public, "
        "mais rarement une telemedecine integree avec dossier chiffre et IA conversationnelle "
        "contextualisee. Au niveau national, les solutions e-sante restent generalistes. "
        "HeraCare se positionne sur le croisement FemTech + telemedecine + IA + conformite 09-08."
    )
    pdf.h2("Besoins fonctionnels")
    for b in [
        "Inscription / connexion (roles : patiente, medecin, admin).",
        "Journal de symptomes et debuts de regles.",
        "Insights predictifs (prochaines regles, ovulation, alertes).",
        "Dossier medical partage chiffre.",
        "Messagerie temps reel securisee.",
        "Hera AI : explication, brief, plan, parse langage naturel.",
    ]:
        pdf.bullet(b)
    pdf.h2("Besoins non-fonctionnels")
    for b in [
        "Securite : JWT, RBAC, AES-256-CBC, bcrypt.",
        "Performance : API REST reactive, Socket.io.",
        "Portabilite : Docker Compose.",
        "Confidentialite : Loi 09-08 / bonnes pratiques HIPAA (chiffrement au repos).",
    ]:
        pdf.bullet(b)

    pdf.add_page()
    pdf.h1("Chapitre 2 - Conception et modelisation")
    pdf.h2("Architecture 3-tiers (MERN)")
    pdf.body(
        "Presentation : React 18, Vite, Tailwind, Recharts. "
        "Logique metier : Node.js, Express, Socket.io, Claude API. "
        "Persistance : MongoDB Atlas (Mongoose)."
    )
    pdf.body(
        "Choix MongoDB : schemas souples pour journaux heterogenes, scalabilite, "
        "integration native Node. Docker garantit la reproductibilite (frontend, backend, mongo)."
    )
    pdf.h2("Schema de donnees")
    pdf.bullet("users : identite, role, assignedDoctor, specialty, password bcrypt")
    pdf.bullet("healthrecords : allergies/medications/notes chiffres AES-256")
    pdf.bullet("symptomlogs : period_start, symptomes, painLevel, mood, flow")
    pdf.bullet("messages : encryptedContent, conversationId")
    pdf.h2("Securite AES-256")
    pdf.body(
        "Cle derivee SHA-256, IV aleatoire 16 octets, stockage iv:ciphertext. "
        "Dechiffrement serveur uniquement apres JWT + controle de role."
    )
    pdf.h2("Engine predictif")
    pdf.body(
        "Moyenne des cycles, prediction J+N, fenetre d'ovulation J-14 +/-2, "
        "alertes : irregularite (ecart-type >= 7 j), douleur recurrente >= 7/10, "
        "signaux SOPK (acne, hirsutisme, prise de poids)."
    )

    pdf.add_page()
    pdf.h1("Chapitre 3 - Realisation")
    pdf.h2("Environnement")
    pdf.body(
        "IDE Cursor/VS Code, Express 4, Mongoose 8, Socket.io 4, React 18, Vite 5, "
        "Tailwind 3, Anthropic Claude (Hera AI), MongoDB Atlas."
    )
    pdf.h2("Modules")
    for m in [
        "Auth - register/login JWT, assignation medecin",
        "Health - symptomes, insights, dossier chiffre",
        "Chat - Socket.io + messages chiffres",
        "AI - /api/ai/* (chat, explain, parse, brief, wellness, ask-doctor)",
        "Crypto - utils/crypto.js AES-256-CBC",
    ]:
        pdf.bullet(m)
    pdf.h2("Hera AI")
    pdf.body(
        "L'assistant contextualise les reponses avec l'historique de cycle et le dossier. "
        "Disclaimer systematique : information generale, pas un diagnostic medical."
    )

    pdf.add_page()
    pdf.h1("Chapitre 4 - Tests et deploiement")
    pdf.h2("Tests")
    pdf.bullet("Parcours manuels patiente / medecin")
    pdf.bullet("Seed MongoDB (scenarios d'alertes)")
    pdf.bullet("Healthcheck /api/healthcheck (encryption + ai)")
    pdf.h2("Deploiement")
    pdf.body(
        "Frontend : Vercel (SPA Vite). Backend : service Node persistant (Render/Railway) "
        "+ MongoDB Atlas. Variables : VITE_API_URL, VITE_SOCKET_URL, JWT_SECRET, "
        "AES_SECRET_KEY, ANTHROPIC_API_KEY, MONGODB_URI, CLIENT_URL."
    )

    pdf.h1("Conclusion et perspectives")
    pdf.body(
        "HeraCare demontre la faisabilite d'une plateforme FemTech securisee integrant "
        "analyse predictive, telemedecine et IA. Perspectives : modeles ML SOPK, "
        "notifications push, teleconsultation video WebRTC, certification ISO 27001 / HDS."
    )

    pdf.ln(10)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(100, 100, 100)
    pdf.multi_cell(
        0, 5,
        T("Document academique - MyHeath PFE 2025-2026 - Nezha Fekoussa"),
    )

    path = OUT / "RAPPORT_PFE_HERACARE.pdf"
    pdf.output(str(path))
    print(f"PDF -> {path}")
    return path


def add_slide(prs, title, bullets):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ROSE
    shape.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ROSE
    p.font.name = "Calibri"

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = b
        para.level = 0
        para.font.size = Pt(18)
        para.font.color.rgb = INK
        para.font.name = "Calibri"
        para.space_after = Pt(10)
    return slide


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xFA, 0xF7, 0xF5)
    bar.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ROSE
    accent.line.fill.background()

    t = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = "HeraCare"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = ROSE
    p2 = tf.add_paragraph()
    p2.text = "Plateforme intelligente de telemedecine & suivi de sante feminine"
    p2.font.size = Pt(22)
    p2.font.color.rgb = INK
    p3 = tf.add_paragraph()
    p3.text = "PFE Ingenieur - MERN / Docker / AES-256 / Socket.io / Claude AI"
    p3.font.size = Pt(16)
    p3.font.color.rgb = MUTED
    p4 = tf.add_paragraph()
    p4.text = "Nezha Fekoussa - 2025/2026"
    p4.font.size = Pt(14)
    p4.font.color.rgb = MUTED

    slides = [
        ("Sommaire", [
            "1. Contexte & problematique",
            "2. Objectifs du projet",
            "3. Architecture MERN & Docker",
            "4. Securite AES-256 & Loi 09-08",
            "5. Algorithme predictif",
            "6. Hera AI (Claude)",
            "7. Telemedecine temps reel",
            "8. Deploiement & perspectives",
        ]),
        ("Contexte FemTech", [
            "Sante feminine encore sous-adressee numeriquement",
            "Deserts medicaux -> besoin de telemedecine",
            "Marche FemTech mondial en forte croissance",
            "Offre nationale fragmentee -> opportunite HeraCare",
        ]),
        ("Problematique", [
            "Proteger les donnees de sante sensibles",
            "Apporter une valeur clinique (prediction + IA)",
            "Architecture modulaire digne d'un PFE ingenieur",
            "Experience utilisateur moderne et rassurante",
        ]),
        ("Objectifs", [
            "Stack MERN + Docker Compose",
            "Chiffrement AES-256 au repos",
            "Engine de cycle & alertes SOPK / endometriose",
            "Hera AI : chat, brief medecin, journal NL",
            "Deploiement Vercel + cloud API + Atlas",
        ]),
        ("Architecture 3-tiers", [
            "Frontend : React / Vite / Tailwind / Recharts (Vercel)",
            "Backend : Node / Express / Socket.io / Claude API",
            "Persistance : MongoDB Atlas (Mongoose)",
            "Modules : Auth / Health / Chat / AI / Crypto",
        ]),
        ("Securite avancee", [
            "JWT + RBAC (patient / doctor / admin)",
            "Mots de passe : bcrypt (cost 12)",
            "AES-256-CBC : dossier medical + messages",
            "Alignement Loi 09-08 / bonnes pratiques HIPAA",
        ]),
        ("Algorithme predictif", [
            "Moyenne historique des cycles",
            "Prochaines regles + fenetre d'ovulation",
            "Alertes : irregularite, douleur severe, signaux SOPK",
            "recommendConsultation -> orientation medecin",
        ]),
        ("Hera AI - wow features", [
            "Chat contextualise sur le suivi patiente",
            "Journal symptomes en langage naturel",
            "Explication narrative des insights",
            "Brief medecin + questions de consultation",
            "Plan bien-etre selon la phase du cycle",
        ]),
        ("Telemedecine", [
            "Socket.io : chat prive Patient <-> Medecin",
            "Contenu chiffre au repos",
            "Assignation medecin depuis le dashboard",
            "Indicateurs de canal securise dans l'UI",
        ]),
        ("Deploiement", [
            "Frontend -> Vercel",
            "Backend -> service Node (Render / Railway)",
            "Base -> MongoDB Atlas",
            "CI : build Vite + variables d'environnement",
        ]),
        ("Conclusion & perspectives", [
            "Plateforme FemTech complete et securisee",
            "IA + prediction + teleconsultation textuelle",
            "Suite : ML SOPK, WebRTC video, notifications, HDS",
            "Merci - Questions ?",
        ]),
    ]

    for title, bullets in slides:
        add_slide(prs, title, bullets)

    path = OUT / "PRESENTATION_PFE_HERACARE.pptx"
    prs.save(str(path))
    print(f"PPTX -> {path}")
    return path


if __name__ == "__main__":
    build_pdf()
    build_pptx()
